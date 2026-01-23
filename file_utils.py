import os
import io
import zipfile

def extract_text_content(file_bytes: bytes, extension: str) -> str | None:
    """
    Best-effort text extraction for common office/PDF formats to render in browser.
    Returns string content or None if extraction fails/not supported.
    """
    try:
        ext = (extension or '').lower()
        
        # Plain Text
        if ext in ['.txt', '.csv', '.json', '.xml', '.md', '.py', '.js', '.html', '.css']:
            try:
                return file_bytes.decode('utf-8')
            except:
                try:
                    return file_bytes.decode('latin-1')
                except:
                    return None

        # PDF
        if ext == '.pdf':
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(file_bytes))
                texts = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t: texts.append(t)
                return '\n\n'.join(texts).strip() or None
            except ImportError:
                return "[Server Error] PyPDF2 not installed."
            except Exception:
                return None

        # DOCX
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                return '\n'.join([p.text for p in doc.paragraphs if p.text]).strip() or None
            except ImportError:
                return "[Server Error] python-docx not installed."
            except Exception:
                return None

        # XLSX
        if ext == '.xlsx':
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
                parts = []
                for ws in wb.worksheets:
                    parts.append(f"--- Sheet: {ws.title} ---")
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(c) if c is not None else '' for c in row]
                        if any(vals): parts.append('\t'.join(vals))
                return '\n'.join(parts).strip() or None
            except ImportError:
                return "[Server Error] openpyxl not installed."
            except Exception:
                return None

        # PPTX
        if ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            texts.append(shape.text)
                return '\n\n'.join(texts).strip() or None
            except ImportError:
                return "[Server Error] python-pptx not installed."
            except Exception:
                return None
                
    except Exception as e:
        print(f"Extraction error: {e}")
        return None
    
    return None

def detect_mime_type(file_bytes: bytes) -> tuple[str, str]:
    """
    Detects extension and MIME type from binary signatures.
    Returns (extension, mime_type).
    """
    try:
        b = file_bytes[:512]
        
        if b.startswith(b"%PDF"): return ('.pdf', 'application/pdf')
        if b.startswith(b"\x89PNG\r\n\x1a\n"): return ('.png', 'image/png')
        if b.startswith(b"\xFF\xD8\xFF"): return ('.jpg', 'image/jpeg')
        if b.startswith(b"GIF8"): return ('.gif', 'image/gif')
        if b.startswith(b"BM"): return ('.bmp', 'image/bmp')
        
        # ZIP / Office Open XML
        if b.startswith(b"PK\x03\x04"):
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    names = set(zf.namelist())
                    if any(n.startswith('word/') for n in names):
                        return ('.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                    if any(n.startswith('xl/') for n in names):
                        return ('.xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                    if any(n.startswith('ppt/') for n in names):
                        return ('.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                return ('.zip', 'application/zip')
            except:
                return ('.zip', 'application/zip')

        # Legacy Office (Compound File Binary)
        if b.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"):
            return ('.doc', 'application/msword') # Generic guess
            
    except Exception:
        pass
        
    return ('', 'application/octet-stream')

def estimate_page_count(file_bytes: bytes, extension: str, text_content: str = None) -> float:
    """
    Estimates the number of pages in the document.
    - PDF: Actual page count.
    - Text/Office: Char count / 3000.
    - Images/Other: Default to 1.
    """
    try:
        ext = (extension or '').lower()
        
        # PDF - Actual Count
        if ext == '.pdf':
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(file_bytes))
                return len(reader.pages)
            except:
                return 1.0
                
        # Text Based - Char Count
        if text_content:
            # Approx 3000 chars per single spaced page
            return max(1.0, len(text_content) / 3000.0)
            
        # Fallback
        return 1.0
        
    except Exception:
        return 1.0
